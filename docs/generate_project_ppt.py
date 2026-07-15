from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = BASE_DIR.parent
OUTPUT_PPT = BASE_DIR / "WT_Automation_项目讲解材料.pptx"
OUTPUT_MD = BASE_DIR / "WT_Automation_项目讲解材料_讲稿.md"

ACCENT = RGBColor(47, 84, 235)
ACCENT_2 = RGBColor(24, 144, 255)
TEXT = RGBColor(31, 35, 40)
MUTED = RGBColor(89, 89, 89)
BG = RGBColor(248, 250, 252)
CARD_BG = RGBColor(255, 255, 255)
OK = RGBColor(56, 142, 60)
WARN = RGBColor(230, 145, 56)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.3), Inches(12.2), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Microsoft YaHei"
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(4)


def add_card(slide, left, top, width, height, title, bullets, title_color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(220, 224, 230)

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = title_color

    content = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.48), width - Inches(0.44), height - Inches(0.62))
    tf = content.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT
        p.space_after = Pt(3)
        first = False


def add_chip(slide, left, top, width, text, fill_rgb):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.36)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill_rgb
    chip.line.color.rgb = fill_rgb
    tf = chip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)


def add_picture_if_exists(slide, path, left, top, width=None, height=None):
    if not path.exists():
        return
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_flow_box(slide, left, top, width, height, text, fill_rgb):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)


def add_arrow_text(slide, left, top, text):
    box = slide.shapes.add_textbox(left, top, Inches(0.55), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = MUTED


def add_kpi(slide, left, top, width, title, value, fill_rgb):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.82)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_rgb
    box.line.color.rgb = fill_rgb

    tf = box.text_frame
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.name = "Microsoft YaHei"
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title
    r2.font.name = "Microsoft YaHei"
    r2.font.size = Pt(9)
    r2.font.color.rgb = RGBColor(255, 255, 255)


def build_versioned_path(path: Path) -> Path:
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return path.with_name(f"{path.stem}_{stamp}{path.suffix}")


def save_with_fallback(prs: Presentation, path: Path) -> Path:
    try:
        prs.save(str(path))
        return path
    except PermissionError:
        fallback = build_versioned_path(path)
        prs.save(str(fallback))
        return fallback


def write_text_with_fallback(path: Path, content: str) -> Path:
    try:
        path.write_text(content, encoding="utf-8")
        return path
    except PermissionError:
        fallback = build_versioned_path(path)
        fallback.write_text(content, encoding="utf-8")
        return fallback


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "WT_Automation 项目讲解", "面向 WT 仿真软件的工程化桌面自动化平台 | 5 页版")
    add_chip(slide, Inches(0.5), Inches(1.15), Inches(1.55), "项目目标", ACCENT)
    add_kpi(slide, Inches(0.5), Inches(1.58), Inches(1.7), "页数控制", "5 页内", ACCENT)
    add_kpi(slide, Inches(2.35), Inches(1.58), Inches(1.7), "执行策略", "四层能力", ACCENT_2)
    add_kpi(slide, Inches(4.2), Inches(1.58), Inches(1.7), "当前定位", "平台化建设", OK)
    add_card(
        slide,
        Inches(0.5),
        Inches(2.6),
        Inches(5.9),
        Inches(4.1),
        "项目背景与痛点",
        [
            "WT 仿真软件界面链路长、窗口类型混合、人工录入参数多，重复劳动明显。",
            "复杂 WPF 弹窗、空标题容器、自绘控件和 Win32 对话框并存，导致单一定位方法不稳定。",
            "纯 AI 全流程可行但成本高、速度慢，对重复步骤存在反复判断和状态回退问题。",
            "因此项目目标不是做一次性脚本，而是建设可维护、可复用、可扩展的自动化平台。",
            "当前已从“流程跑通”升级为“流程定义、控件资产、兜底恢复、结果回溯”闭环建设。",
        ],
    )
    add_card(
        slide,
        Inches(6.7),
        Inches(2.6),
        Inches(6.1),
        Inches(4.1),
        "本阶段建设目标",
        [
            "以 Robot Framework 为调度入口，以 Python 为执行核心，建立配置化的 WT 自动化主链路。",
            "把固定步骤沉淀为结构化 action/flow_ref，而不是继续堆积零散录制脚本。",
            "形成“结构化定位 + 相对区域交互 + 模板匹配 + AI 介入”的多层执行能力。",
            "建设总控台、流程编辑器、控件库和模板库，降低后续流程维护成本。",
            "通过运行报告、调试截图和最小回归测试提高问题定位效率与交付稳定性。",
            "将项目逐步沉淀为可复用的桌面自动化解决方案，而非单一业务 Demo。",
        ],
        title_color=ACCENT_2,
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "技术路线与关键机制", "从全流程 AI 试探转向以结构化执行为主、分层兜底为辅的工程化路线")
    add_flow_box(slide, Inches(0.55), Inches(2.2), Inches(2.25), Inches(1.0), "阶段 1\n全流程 AI", WARN)
    add_arrow_text(slide, Inches(2.95), Inches(2.45), "→")
    add_flow_box(slide, Inches(3.45), Inches(2.2), Inches(2.55), Inches(1.0), "阶段 2\n结构化动作主执行", ACCENT)
    add_arrow_text(slide, Inches(6.15), Inches(2.45), "→")
    add_flow_box(slide, Inches(6.65), Inches(2.2), Inches(2.55), Inches(1.0), "阶段 3\n模板/相对区域补偿", ACCENT_2)
    add_arrow_text(slide, Inches(9.35), Inches(2.45), "→")
    add_flow_box(slide, Inches(9.85), Inches(2.2), Inches(2.55), Inches(1.0), "阶段 4\nAI 续跑恢复", OK)
    add_card(
        slide,
        Inches(0.7),
        Inches(3.7),
        Inches(5.9),
        Inches(2.55),
        "路线调整原因",
        [
            "固定步骤由结构化动作执行，速度更快、可解释性更强，适合稳定业务链路。",
            "WPF 弹窗、自绘按钮、离屏文本层等场景难以靠纯控件定位稳定命中，需要相对区域和模板补偿。",
            "AI 不再从头操作全流程，而是仅处理当前卡住步骤，显著降低 token 成本和重复执行风险。",
            "主路径保持可验证，异常路径再交给模板和 AI，整体稳定性与运维成本更平衡。",
        ],
    )
    add_card(
        slide,
        Inches(6.8),
        Inches(3.7),
        Inches(5.8),
        Inches(2.55),
        "关键机制",
        [
            "控件定位采用多属性评分：综合 name、automation_id、class_name、framework_id 等信息筛选候选。",
            "复杂窗口优先锁定父窗口，再按相对区域换算点击/输入坐标，提升跨分辨率和弱控件场景适应性。",
            "首次点击导致前台窗口切换时，重选更具体的参考窗口并补点，修正坐标偏移。",
            "记录 `resume_stage / fallback_stage` 和失败上下文，让 AI 从指定阶段续跑而非重做全链路。",
        ],
        title_color=OK,
    )
    add_chip(slide, Inches(0.7), Inches(6.45), Inches(3.0), "亮点：AI 从重执行改为续跑", OK)
    add_chip(slide, Inches(3.95), Inches(6.45), Inches(3.45), "亮点：相对区域解决复杂 WPF 场景", ACCENT_2)
    add_chip(slide, Inches(7.7), Inches(6.45), Inches(4.1), "亮点：恢复链路支持模板与 AI 双兜底", ACCENT)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "总体架构与项目分层", "主调度、流程定义、执行定位、资产沉淀与结果回溯已经形成完整链路")
    add_kpi(slide, Inches(0.55), Inches(1.0), Inches(1.7), "调度入口", "Robot", ACCENT)
    add_kpi(slide, Inches(2.45), Inches(1.0), Inches(1.7), "执行核心", "Python", ACCENT_2)
    add_kpi(slide, Inches(4.35), Inches(1.0), Inches(1.7), "定义层", "JSON/Excel", WARN)
    add_kpi(slide, Inches(6.25), Inches(1.0), Inches(1.7), "智能介入", "UI-TARS", OK)
    add_kpi(slide, Inches(8.15), Inches(1.0), Inches(1.7), "维护入口", "WT_Launcher", ACCENT)
    add_kpi(slide, Inches(10.05), Inches(1.0), Inches(1.7), "定义入口", "WT_Flow_Editor", ACCENT_2)
    add_flow_box(slide, Inches(0.5), Inches(1.7), Inches(2.15), Inches(0.9), "Robot Framework\n总调度入口", ACCENT)
    add_arrow_text(slide, Inches(2.75), Inches(1.95), "→")
    add_flow_box(slide, Inches(3.2), Inches(1.7), Inches(2.2), Inches(0.9), "Python 主流程\nWT_AUT_recorded.py", ACCENT_2)
    add_arrow_text(slide, Inches(5.55), Inches(1.95), "→")
    add_flow_box(slide, Inches(6.0), Inches(1.7), Inches(1.9), Inches(0.9), "执行器 /\n定位器", OK)
    add_arrow_text(slide, Inches(8.0), Inches(1.95), "→")
    add_flow_box(slide, Inches(8.45), Inches(1.7), Inches(1.9), Inches(0.9), "模板/控件库\n资产层", WARN)
    add_arrow_text(slide, Inches(10.45), Inches(1.95), "→")
    add_flow_box(slide, Inches(10.9), Inches(1.7), Inches(1.9), Inches(0.9), "UI-TARS\nAI 介入", RGBColor(84, 110, 122))
    add_card(
        slide,
        Inches(0.55),
        Inches(3.0),
        Inches(4.05),
        Inches(2.95),
        "核心执行层",
        [
            "主入口 `WT_AUT_recorded.py` 负责读取配置、组装运行参数并调度流程。",
            "执行器 `wt_flow_executor.py` 统一解释 click、type_text、click_relative_region 等动作。",
            "定位器 `wt_flow_locator.py` 负责窗口筛选、控件评分、相对区域参考窗口判断与补点。",
            "业务辅助模块承接投影、窗口激活、文件对话框、业务步骤等专用逻辑。",
        ],
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(3.0),
        Inches(3.75),
        Inches(2.95),
        "定义与资产层",
        [
            "`flow_definition.json` 与 `flow_packages/` 负责流程编排和步骤组织。",
            "`control_maps/` 按窗口和框架沉淀控件资产，支持从编辑器回写与复用。",
            "`image_templates/` 保存模板图片，用于特殊按钮、图标和弹层的视觉匹配。",
            "录制结果可通过转换器提升为正式流程定义，并提取运行参数占位。",
        ],
        title_color=ACCENT_2,
    )
    add_card(
        slide,
        Inches(8.75),
        Inches(3.0),
        Inches(4.05),
        Inches(2.95),
        "维护与质量层",
        [
            "`WT_Launcher.py` 提供总控台，支持运行、日志查看、模型配置和工具入口整合。",
            "`WT_Flow_Editor.py` 提供步骤模板、新增相对区域、控件搜索和控件库导入。",
            "`wt_run_reporting.py` 输出结构化运行报告，记录步骤结果、耗时、fallback 与失败原因。",
            "`tests/` 中已建立最小回归测试，覆盖执行器、编辑器工具、运行报告和转换器。",
        ],
        title_color=OK,
    )

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "关键能力与使用入口", "围绕日常运行、异常恢复、流程维护和资产沉淀提供可视化支撑")
    add_chip(slide, Inches(0.5), Inches(1.0), Inches(2.25), "使用对象：运行人员 / 调试人员 / 维护人员", ACCENT)
    add_chip(slide, Inches(2.95), Inches(1.0), Inches(2.2), "使用方式：总控台统一进入", ACCENT_2)
    add_card(
        slide,
        Inches(0.5),
        Inches(1.35),
        Inches(4.0),
        Inches(2.7),
        "1. WT_Launcher 总控台",
        [
            "用于启动/停止主流程，实时查看日志、步骤状态和结构化运行报告。",
            "支持模型配置持久化、最近历史管理、环境检查和常用工具统一入口。",
            "运行报告界面可快速查看每个步骤的成功/失败、耗时、执行策略和错误信息。",
            "适合作为项目演示、运维值守和问题排查的统一入口。",
        ],
    )
    add_card(
        slide,
        Inches(0.5),
        Inches(4.0),
        Inches(4.0),
        Inches(2.35),
        "2. 模板采集与控件资产",
        [
            "模板采集器用于制作按钮、图标和特殊弹层模板，补足结构化定位盲区。",
            "控件信息可自动沉淀到 `control_maps/`，按窗口标题和框架类型分类保存。",
            "模板库和控件库共同承担资产层角色，既服务执行，也服务后续编辑与复用。",
            "建议模板采集与运行环境保持一致的缩放比例，以提高命中稳定性。",
        ],
        title_color=WARN,
    )
    add_card(
        slide,
        Inches(4.75),
        Inches(1.35),
        Inches(4.0),
        Inches(5.0),
        "3. WT_Flow_Editor 流程编辑器",
        [
            "支持步骤模板化新增，已预置按钮、输入框、下拉项和相对区域等高频模板。",
            "可维护 action 字段、控件信息、辅助判断、fallback 链路和运行参数占位。",
            "支持从 Inspect 文本解析、控件库导入、半自动采集和交互点击采集获取控件信息。",
            "新增父窗口相对区域动作配置，便于维护 WPF 弹窗、自绘输入框和难命中列表项。",
            "这是将零散脚本升级为可配置流程的关键入口，也是后续产品化的基础。",
        ],
        title_color=ACCENT_2,
    )
    add_card(
        slide,
        Inches(9.0),
        Inches(1.35),
        Inches(3.85),
        Inches(5.0),
        "4. 关键机制概览",
        [
            "主路径：优先走结构化控件定位和 action 执行。",
            "复杂 WPF 场景：锁定父窗口后按相对区域换算真实点击/输入位置。",
            "首次点击导致窗口切换时：重选更具体的参考窗口并执行补点，修正坐标偏移。",
            "模板兜底：结构化定位失败后，用模板匹配恢复局部动作。",
            "AI 介入：模板也失败时，仅针对当前卡住步骤续跑，不重做全流程。",
            "建议使用顺序：先总控台运行，再模板补强，最后在编辑器维护流程与控件资产。",
        ],
        title_color=OK,
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "阶段成果、价值与下一步", "项目已具备工程化雏形，后续重点转向能力联动、数据沉淀与稳定性深化")
    add_kpi(slide, Inches(0.5), Inches(1.0), Inches(1.85), "路线状态", "已成型", OK)
    add_kpi(slide, Inches(2.55), Inches(1.0), Inches(1.85), "主控台", "已落地", ACCENT)
    add_kpi(slide, Inches(4.6), Inches(1.0), Inches(1.85), "流程编辑", "已可用", ACCENT_2)
    add_kpi(slide, Inches(6.65), Inches(1.0), Inches(1.85), "分层兜底", "已验证", WARN)
    add_kpi(slide, Inches(8.7), Inches(1.0), Inches(1.85), "运行报告", "已接入", OK)
    add_kpi(slide, Inches(10.75), Inches(1.0), Inches(1.85), "后续重点", "平台联动", ACCENT)
    add_card(
        slide,
        Inches(0.5),
        Inches(1.45),
        Inches(4.0),
        Inches(4.9),
        "当前成果",
        [
            "WT 主流程已能承接多类典型业务节点，技术路线已从试验阶段进入稳定打磨阶段。",
            "总控台、流程编辑器、模板采集器、运行报告和最小回归测试已形成配套工具链。",
            "流程定义、控件库和模板库开始持续沉淀，降低后续新链路接入和维护成本。",
            "项目已不再依赖单一脚本堆叠，而是形成“执行层 + 定义层 + 资产层 + 维护层”的平台雏形。",
        ],
        title_color=ACCENT,
    )
    add_card(
        slide,
        Inches(4.7),
        Inches(1.45),
        Inches(4.0),
        Inches(4.9),
        "项目价值",
        [
            "相比纯人工操作，可显著减少重复录入和重复点击，提高业务执行效率。",
            "相比纯录制脚本，当前系统更适合处理 WPF/Win32 混合场景，具备更强可恢复能力。",
            "相比纯 AI 执行，现方案保留了稳定、低成本和可解释的主路径，更适合长期交付。",
            "结构化运行报告和步骤级问题回溯提升了调试效率，也为论文实验和后续考核提供了数据基础。",
        ],
        title_color=WARN,
    )
    add_card(
        slide,
        Inches(8.9),
        Inches(1.45),
        Inches(3.9),
        Inches(4.9),
        "下一步计划",
        [
            "继续联通流程编辑器定义与实际执行，让配置修改可以更快回流到运行层。",
            "完善 AI 介入的实验结果沉淀，把模板兜底与 AI 续跑效果做量化对比。",
            "继续补齐覆盖区、格网、导出等尾部链路的稳定性验证和报告统计。",
            "优化复杂控件采集与定位，进一步提升特定菜单、小图标和树节点命中率。",
            "把项目沉淀为可迁移的桌面自动化方案模板，服务更多复杂 GUI 软件场景。",
        ],
        title_color=OK,
    )

    return save_with_fallback(prs, OUTPUT_PPT)


def build_notes():
    content = """# WT_Automation 项目讲解讲稿

## 第1页 项目背景与目标
- 这一页先把项目定位讲清楚：WT_Automation 不是单一脚本，而是面向 WT 仿真软件的桌面自动化平台。
- 之所以要做这个项目，是因为 WT 界面流程长、参数录入多，而且混合了 WPF 和 Win32 两类窗口，人工操作重复度高、出错率也高。
- 项目目标不是简单追求 AI 全自动，而是建设稳定、可维护、可复用的执行体系。
- 当前已经从“先跑通”走到了“可配置、可兜底、可回溯、可迭代”的阶段。

## 第2页 技术路线与关键机制
- 项目最早尝试过全流程 AI，但很快发现成本高、速度慢，而且重复步骤会被反复识别。
- 现在的核心路线是：结构化动作主执行，复杂场景用相对区域和模板补偿，实在恢复不了再由 AI 介入。
- 这里有两个关键点要强调：第一，复杂 WPF 场景通过父窗口相对区域解决弱控件定位问题；第二，AI 只处理当前卡住步骤，而不是从头重跑。
- 另外，首次点击如果导致前台窗口切换，系统会重选更具体的参考窗口并补点，这对弹窗场景很关键。

## 第3页 总体架构与项目分层
- 架构上，Robot Framework 负责总调度，Python 负责执行编排和能力集成。
- 执行层里有 `wt_flow_executor.py` 和 `wt_flow_locator.py`，分别负责动作执行和控件定位。
- 定义层有 `flow_definition.json`、`flow_packages/`、流程编辑器和录制转换器，支撑流程维护。
- 资产层有 `control_maps/` 和 `image_templates/`，支撑控件复用和模板兜底。
- 维护层则通过总控台、运行报告、调试截图和最小回归测试构成闭环。

## 第4页 关键能力与使用入口
- 日常运行入口是 `WT_Launcher` 总控台，可以运行流程、看日志、看步骤结果和打开工具。
- 当结构化定位不稳定时，可以用模板采集器补模板，也可以让控件信息自动沉淀到控件库。
- 当要新增或维护流程时，进入 `WT_Flow_Editor`，通过模板化步骤、控件搜索和相对区域配置快速编辑。
- 这一页要重点说明：项目已经不靠纯代码手改，而是逐步形成了面向运行、维护和调试的可视化工具链。

## 第5页 阶段成果、价值与下一步
- 当前项目已经具备工程化雏形，主链路、编辑器、模板采集器、运行报告和回归测试都已经落地。
- 对业务而言，它降低了人工操作成本；对技术而言，它比纯录制脚本更稳，比纯 AI 执行更可控。
- 下一步重点是把流程编辑器和实际执行进一步联通，并补齐更多尾部链路的稳定性验证。
- 同时也要持续优化复杂控件的命中率，并沉淀更多实验数据，让项目既能跑、也能讲、还能写进论文。
"""
    return write_text_with_fallback(OUTPUT_MD, content)


if __name__ == "__main__":
    ppt_path = build_presentation()
    note_path = build_notes()
    print(f"generated: {ppt_path}")
    print(f"generated: {note_path}")
